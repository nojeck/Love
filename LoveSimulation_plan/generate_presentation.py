import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False
import requests


def read_md(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def sections_from_md(text):
    lines = text.splitlines()
    sections = []
    current = None
    for line in lines:
        if line.strip().startswith('#'):
            title = line.strip().lstrip('#').strip()
            current = {"title": title, "content": []}
            sections.append(current)
        else:
            if current is None:
                current = {"title": "Intro", "content": []}
                sections.append(current)
            current["content"].append(line)
    for s in sections:
        s["content"] = " ".join([l.strip() for l in s["content"] if l.strip()])
    return sections


def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0


def add_two_col_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    # left column
    p = tf.add_paragraph()
    p.text = left_title
    p.level = 0
    for b in left_bullets:
        p = tf.add_paragraph()
        p.text = f"- {b}"
        p.level = 1
    # spacer
    p = tf.add_paragraph()
    p.text = ""
    # right column (as more bullets)
    p = tf.add_paragraph()
    p.text = right_title
    p.level = 0
    for b in right_bullets:
        p = tf.add_paragraph()
        p.text = f"- {b}"
        p.level = 1


def ensure_sample_images(folder):
    if not PIL_AVAILABLE:
        return []
    os.makedirs(folder, exist_ok=True)
    paths = []
    for name, color in [("cover", (70,130,180)), ("episode1", (200,80,80))]:
        p = os.path.join(folder, f"{name}.png")
        if not os.path.exists(p):
            img = Image.new('RGB', (1200, 600), color)
            d = ImageDraw.Draw(img)
            try:
                f = ImageFont.truetype("malgun.ttf", 48)
            except Exception:
                f = ImageFont.load_default()
            d.text((40, 40), f"{name} image", font=f, fill=(255,255,255))
            img.save(p)
        paths.append(p)
    return paths


def download_images(urls, folder):
    os.makedirs(folder, exist_ok=True)
    paths = []
    for i, url in enumerate(urls):
        ext = 'jpg'
        p = os.path.join(folder, f'download_{i}.{ext}')
        try:
            r = requests.get(url, timeout=15)
            if r.status_code == 200:
                with open(p, 'wb') as f:
                    f.write(r.content)
                paths.append(p)
        except Exception:
            continue
    return paths


def set_korean_font_on_textframe(tf, font_name='Malgun Gothic', size_pt=18):
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size_pt)
        if not paragraph.runs:
            r = paragraph.add_run()
            r.text = ''
            r.font.name = font_name
            r.font.size = Pt(size_pt)


def add_image_slide(prs, title, image_path, caption=None):
    slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    if title:
        try:
            slide.shapes.title.text = title
            set_korean_font_on_textframe(slide.shapes.title.text_frame)
        except Exception:
            pass
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width)
    if caption:
        tx = slide.shapes.add_textbox(left, top+Inches(4.8), width, Inches(0.8)).text_frame
        tx.text = caption
        set_korean_font_on_textframe(tx)


def add_branching_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # root node
    left = Inches(3)
    top = Inches(0.8)
    w = Inches(4)
    h = Inches(0.8)
    root = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    root.text_frame.text = "Question: '나 오늘 뭐 달라진 거 없어?'"
    set_korean_font_on_textframe(root.text_frame, size_pt=14)
    root.fill.solid()
    root.fill.fore_color.rgb = RGBColor(70,130,180)
    root.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)

    # left child - correct
    l_left = Inches(1)
    l_top = Inches(2.4)
    l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_left, l_top, Inches(3.5), Inches(0.8))
    l.text_frame.text = "정답: 귀걸이 언급\n-> 통과"
    set_korean_font_on_textframe(l.text_frame, size_pt=12)
    l.fill.solid(); l.fill.fore_color.rgb = RGBColor(102,205,170)

    # right child - wrong
    r_left = Inches(5)
    r_top = Inches(2.4)
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, r_left, r_top, Inches(3.5), Inches(0.8))
    r.text_frame.text = "오답: 부정적 반응\n-> 병맛 회귀"
    set_korean_font_on_textframe(r.text_frame, size_pt=12)
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(240,128,128)

    # lines
    slide.shapes.add_connector(1, left+Inches(0.2), top+Inches(0.8), l_left+Inches(0.5), l_top)
    slide.shapes.add_connector(1, left+Inches(3.8), top+Inches(0.8), r_left+Inches(0.5), r_top)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", "-o", default="LoveSimulation.pptx")
    parser.add_argument("--files", nargs='*', default=[
        "AI 연애 시뮬레이션 기획서 유사 콘텐츠 조사.md",
        "무한 회귀 연애!.md",
    ])
    parser.add_argument("--download-images", action='store_true', help='Download web images to include')
    parser.add_argument("--image-urls", nargs='*', default=[
        "https://picsum.photos/1200/600?random=1",
        "https://picsum.photos/1200/600?random=2",
    ])
    args = parser.parse_args()

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "무한 회귀 연애 — 기획서 요약"
    try:
        subtitle = title_slide.placeholders[1]
        subtitle.text = "생성: 자동화 PPT\n출처: 프로젝트 문서 병합"
    except Exception:
        pass

    bullets = []
    for fp in args.files:
        if not os.path.exists(fp):
            continue
        text = read_md(fp)
        secs = sections_from_md(text)
        if secs:
            first = secs[0]
            excerpt = first["content"][:250] if first["content"] else ""
            bullets.append(f"{first['title']}: {excerpt}")
            for s in secs[1:6]:
                if s["title"] and s["content"]:
                    bullets.append(f"{s['title']}: {s['content'][:200]}")

    for i in range(0, len(bullets), 6):
        add_bullet_slide(prs, "문서 요약", bullets[i:i+6])

    add_bullet_slide(prs, "핵심 시스템", [
        "V.E.A. (Voice Emotion Analysis): STT + Jitter/Shimmer/Pitch/HNR",
        "혼돈 지수 (Chaos Meter): 딴소리/비상식 행위 트래킹",
        "기억 시스템: 이전 루프 기억으로 난이도 상승",
        "설문 데이터 기반 판정: RAG 방식 가중치",
    ])

    add_bullet_slide(prs, "기술 스택 및 구현 가이드", [
        "엔진: Unity",
        "AI: LLM + RAG, Whisper or Native STT",
        "오디오 분석: Librosa / Praat",
        "TTS: ElevenLabs 등",
    ])

    add_two_col_slide(prs, "리스크 및 윤리", 
                      "리스크", 
                      ["음성 판정 오탐: 품질/잡음 보정 필요", "실제 데이터 사용 시 법적·윤리적 문제"],
                      "대응책",
                      ["캘리브레이션/노이즈 필터링", "명확한 동의 화면·익명화·보관정책"])

    add_bullet_slide(prs, "다음 작업(권장)", [
        "Episode 01 프로토타입 개발: STT+키워드 매칭 + 간단 오디오 피처",
        "SER 모델 검증(데이터셋 확보 및 라벨링)",
        "동의 화면 및 개인정보 정책 초안 작성",
        "유저 테스팅(소규모 플레이테스트, 퀄리티 체크)",
    ])

    add_bullet_slide(prs, "Episode 01 - 분기 플로우 (예시)", [
        "질문: '나 오늘 뭐 달라진 거 없어?' (정답: 귀걸이 언급)",
        "정답: 호감도 상승 -> 다음 스테이지",
        "오답: 부정적 반응 -> 병맛 사망 연출(헐크 변신 등)",
        "딴소리/무응답: 혼돈 지수 상승 -> 히든 루트/사망 회귀",
    ])

    add_bullet_slide(prs, "판정 로직 초안 (예시)", [
        "텍스트 매칭: 키워드(귀걸이) 가중치 0.4",
        "오디오 지표: Jitter/ Shimmer/ Pitch/ HNR 합산 가중치 0.4",
        "문맥 신뢰도: STT confidence 0.2",
        "임계치 예: 합계 >= 0.7 -> 통과, <0.7 -> 회귀",
    ])

    add_bullet_slide(prs, "개인정보 동의(샘플)", [
        "녹음 및 음성 데이터 수집에 대한 명시적 동의 필요",
        "익명화, 보관 기간, 이용 목적을 명시",
        "사용자는 언제든 데이터 삭제 요청 가능",
    ])

    # prepare images: either download web images (user permitted) or generate samples
    image_folder = os.path.join(os.path.dirname(__file__), 'images')
    images = []
    if args.download_images:
        images = download_images(args.image_urls, image_folder)
    if not images:
        images = ensure_sample_images(image_folder)

    if images:
        add_image_slide(prs, "Cover", images[0], caption="무한 회귀 연애 - 프로젝트 커버")
        if len(images) > 1:
            add_image_slide(prs, "Episode 01", images[1], caption="Episode 01 예시 이미지")

    add_branching_slide(prs)

    prs.save(args.output)
    print("Saved:", args.output)


if __name__ == '__main__':
    main()
